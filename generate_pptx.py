"""
Générateur de présentation PowerPoint — DÉGG JOJ Dakar 2026
Orange Hackathon Pitch Deck
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# ─── Couleurs DÉGG ───────────────────────────────────────────────────────────
GREEN        = RGBColor(0x0F, 0xA9, 0x58)   # #0FA958 vert principal
GREEN_DARK   = RGBColor(0x0C, 0x8A, 0x48)   # #0C8A48 vert foncé
GREEN_DEEP   = RGBColor(0x0A, 0x2A, 0x16)   # #0A2A16 vert très foncé
GREEN_LIGHT  = RGBColor(0xE8, 0xF8, 0xEF)   # fond vert très clair
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT   = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_MED     = RGBColor(0x9E, 0x9E, 0x9E)
DARK_TEXT    = RGBColor(0x0B, 0x12, 0x20)   # #0B1220
ORANGE_ACC   = RGBColor(0xFF, 0x7A, 0x1A)   # orange Orange Sénégal
ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xE6)

# ─── Dimensions 16:9 ─────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # totalement vide

# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, color, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        5,  # rounded rectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=20, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=16, color=DARK_TEXT, spacing=1.15):
    """lines = [(text, bold, color_override)]"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold, col) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = col if col else color
    return txBox


def add_badge(slide, text, left, top, width=1.4, height=0.38,
              bg=GREEN, fg=WHITE, font_size=13):
    r = add_rounded_rect(slide, left, top, width, height, bg)
    add_text(slide, text, left+0.05, top+0.03, width-0.1, height-0.06,
             font_size=font_size, bold=True, color=fg, align=PP_ALIGN.CENTER)
    return r


def slide_number(slide, n, total=7):
    add_text(slide, f"{n} / {total}",
             12.3, 7.1, 0.9, 0.3,
             font_size=11, color=GRAY_MED, align=PP_ALIGN.RIGHT)


def green_header_bar(slide, title, subtitle=""):
    """Barre verte en haut pour les slides intérieures"""
    add_rect(slide, 0, 0, 13.33, 1.25, GREEN_DEEP)
    add_text(slide, title, 0.4, 0.1, 10, 0.6,
             font_size=30, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.72, 10, 0.42,
                 font_size=15, color=RGBColor(0xA8, 0xE6, 0xC3))


def footer_band(slide, text="DÉGG · JOJ Dakar 2026 · Orange Hackathon"):
    add_rect(slide, 0, 7.2, 13.33, 0.3, GREEN_DEEP)
    add_text(slide, text, 0.3, 7.22, 12, 0.26,
             font_size=10, color=RGBColor(0xA8, 0xE6, 0xC3))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)

# Fond vert foncé sur toute la slide
add_rect(s1, 0, 0, 13.33, 7.5, GREEN_DEEP)

# Motif géométrique décoratif (cercles concentriques)
for r_inch in [1.2, 1.8, 2.4, 3.0]:
    sh = s1.shapes.add_shape(9, Inches(10.5 - r_inch), Inches(-r_inch),
                              Inches(r_inch*2), Inches(r_inch*2))
    sh.fill.background()
    sh.line.color.rgb = RGBColor(0x0F, 0xA9, 0x58)
    sh.line.width = Pt(1.5)

# Bande verte en bas
add_rect(s1, 0, 6.4, 13.33, 1.1, GREEN)

# Logo "D" carré arrondi
logo_bg = add_rounded_rect(s1, 0.6, 1.2, 1.5, 1.5, GREEN)
add_text(s1, "D", 0.6, 1.25, 1.5, 1.4,
         font_size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Nom app
add_text(s1, "DÉGG", 2.4, 1.15, 6, 1.0,
         font_size=64, bold=True, color=WHITE)

# Tagline JOJ
add_text(s1, "JEUX OLYMPIQUES DE LA JEUNESSE", 2.4, 2.15, 8, 0.5,
         font_size=14, bold=True, color=GREEN,
         align=PP_ALIGN.LEFT)

# Trait séparateur
sep = s1.shapes.add_shape(1, Inches(2.4), Inches(2.7), Inches(6.5), Pt(2))
sep.fill.solid()
sep.fill.fore_color.rgb = GREEN
sep.line.fill.background()

# Proposition de valeur
add_text(s1,
         "L'application officielle de traduction & mobilité\npour les JOJ Dakar 2026 · 13 langues dont le Wolof",
         2.4, 2.8, 9.5, 1.0,
         font_size=18, color=RGBColor(0xCC, 0xF0, 0xDD))

# Badges features
badges = [
    ("13 Langues", 0.6, 4.4),
    ("Hors-ligne", 2.4, 4.4),
    ("PWA + Mobile", 4.2, 4.4),
    ("IA Multilingue", 6.0, 4.4),
    ("Carte des sites", 7.9, 4.4),
]
for (lbl, lx, ly) in badges:
    add_badge(s1, lbl, lx, ly, 1.5, 0.4, bg=GREEN, fg=WHITE, font_size=12)

# Dakar 2026 badge
add_badge(s1, "DAKAR 2026", 0.6, 5.1, 2.0, 0.45,
          bg=WHITE, fg=GREEN_DEEP, font_size=14)
add_badge(s1, "Orange Hackathon", 2.9, 5.1, 2.5, 0.45,
          bg=ORANGE_ACC, fg=WHITE, font_size=14)

# Bas de slide — slogan Wolof
add_text(s1, "« Dafa dégg » — Il comprend", 0.5, 6.5, 8, 0.5,
         font_size=16, italic=True, color=GREEN_DEEP, bold=True)
add_text(s1, "Langue nationale : Wolof 🇸🇳", 8.5, 6.5, 4.5, 0.5,
         font_size=13, color=GREEN_DEEP, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CONSTAT
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, 13.33, 7.5, WHITE)
green_header_bar(s2, "Constat", "Le problème que nous résolvons")
footer_band(s2)
slide_number(s2, 2)

# Fond gris clair contenu
add_rect(s2, 0.3, 1.4, 12.73, 5.6, GRAY_LIGHT)

# 4 cartes stats (2x2)
stat_cards = [
    ("206", "Nations représentées aux JOJ", GREEN),
    ("4 000+", "Athlètes de 15 à 18 ans", GREEN_DARK),
    ("50 000+", "Visiteurs internationaux attendus", GREEN),
    ("80 %", "Ne parlent ni français ni wolof", RGBColor(0xC0, 0x39, 0x2B)),
]

positions = [(0.5, 1.6), (3.5, 1.6), (6.5, 1.6), (9.5, 1.6)]
for i, ((num, label, col), (lx, ly)) in enumerate(zip(stat_cards, positions)):
    add_rect(s2, lx, ly, 2.8, 2.0, WHITE)
    add_text(s2, num, lx+0.1, ly+0.08, 2.6, 0.9,
             font_size=40, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s2, label, lx+0.08, ly+0.9, 2.64, 0.9,
             font_size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# Titre section problèmes
add_rect(s2, 0.3, 3.75, 12.73, 0.35, GREEN)
add_text(s2, "LES 3 PROBLÈMES IDENTIFIÉS", 0.5, 3.78, 12, 0.3,
         font_size=13, bold=True, color=WHITE)

# 3 problèmes
problems = [
    ("🚫", "Barrière linguistique",
     "Les bénévoles et visiteurs ne peuvent pas communiquer\nentre 206 nationalités différentes"),
    ("🗺️", "Navigation impossible",
     "8 sites JOJ répartis en 3 zones sans guide\nde transport adapté en langue locale"),
    ("📱", "Aucune solution locale",
     "Les apps génériques (Google Translate) ne connaissent\nni le Wolof, ni les JOJ, ni Dakar"),
]

for i, (icon, title, desc) in enumerate(problems):
    lx = 0.5 + i * 4.2
    add_rect(s2, lx, 4.2, 3.8, 2.7, WHITE)
    add_rect(s2, lx, 4.2, 3.8, 0.5, GREEN)
    add_text(s2, f"  {title}", lx+0.1, 4.22, 3.6, 0.46,
             font_size=14, bold=True, color=WHITE)
    add_text(s2, desc, lx+0.12, 4.78, 3.56, 2.0,
             font_size=13, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, 13.33, 7.5, WHITE)
green_header_bar(s3, "Solution", "Notre proposition de valeur")
footer_band(s3)
slide_number(s3, 3)

# Citation centrale
add_rect(s3, 0.4, 1.4, 12.53, 1.3, GREEN_DEEP)
add_text(s3,
         "DÉGG élimine la barrière linguistique aux JOJ en 13 langues,",
         0.6, 1.48, 12.1, 0.55,
         font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s3,
         "dont le Wolof — sans installation — dans la langue de l'utilisateur",
         0.6, 1.96, 12.1, 0.55,
         font_size=18, color=GREEN, align=PP_ALIGN.CENTER)

# Parcours utilisateur visuel
steps = [
    ("1", "Atterrit\nà Dakar", ""),
    ("2", "Ouvre\nDÉGG", ""),
    ("3", "Choisit\nsa langue", ""),
    ("4", "Utilise\nl'app", ""),
    ("5", "Communique\nlocalement", ""),
]

# Ligne connectrice
add_rect(s3, 0.8, 3.55, 11.73, 0.08, GREEN)

for i, (num, label, _) in enumerate(steps):
    lx = 0.6 + i * 2.4
    # Cercle
    circle = s3.shapes.add_shape(9, Inches(lx), Inches(3.05),
                                  Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    add_text(s3, num, lx, 3.08, 0.7, 0.65,
             font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s3, label, lx-0.3, 3.85, 1.3, 0.7,
             font_size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# Flèches entre cercles
for i in range(4):
    lx = 1.38 + i * 2.4
    add_text(s3, "→", lx, 3.12, 0.5, 0.5,
             font_size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# 5 avantages différenciants
avantages = [
    ("Seule app supportant le Wolof"),
    ("Base de connaissances JOJ unique"),
    ("Interface 100 % dans votre langue"),
    ("Fonctionne hors-ligne (PWA)"),
    ("IA multilingue 24h/24"),
]

add_rect(s3, 0.4, 4.75, 5.8, 2.45, GREEN_LIGHT)
add_text(s3, "DIFFÉRENCIATEURS CLÉS", 0.55, 4.82, 5.5, 0.35,
         font_size=12, bold=True, color=GREEN_DARK)
for i, av in enumerate(avantages):
    add_text(s3, f"✓  {av}", 0.55, 5.22 + i*0.38, 5.5, 0.35,
             font_size=14, color=DARK_TEXT)

# Visuel app simulé (mockup simplifié)
add_rect(s3, 7.0, 4.7, 5.9, 2.6, GREEN_DEEP)
add_rect(s3, 7.15, 4.82, 5.6, 0.35, GREEN)
add_text(s3, "DÉGG", 7.15, 4.84, 2, 0.3,
         font_size=14, bold=True, color=WHITE)
add_text(s3, "🇸🇳 Wolof  →  🇫🇷 Français", 7.15, 5.25, 5.5, 0.35,
         font_size=13, color=RGBColor(0xA8, 0xE6, 0xC3))
add_rect(s3, 7.15, 5.65, 5.6, 0.65, WHITE)
add_text(s3, "Où est le stade ?", 7.25, 5.72, 5.4, 0.5,
         font_size=14, color=DARK_TEXT)
add_rect(s3, 7.15, 6.38, 5.6, 0.65, GREEN)
add_text(s3, "Fan bu toog mi ci dëkk bi ?", 7.25, 6.45, 5.4, 0.5,
         font_size=14, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PRODUIT / SERVICE
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, 13.33, 7.5, WHITE)
green_header_bar(s4, "Produit / Service", "5 fonctionnalités · 1 application")
footer_band(s4)
slide_number(s4, 4)

features = [
    ("01", "Traduction\nInstantanée",
     "Texte + voix · 13 langues\nDébouncé 800ms · Historique",
     "Tapez ou parlez\n→ Traduction en temps réel\n→ Lecture audio"),
    ("02", "Phrases JOJ\nPré-chargées",
     "6 catégories · 48 phrases\nOrientation · Urgences · Sport",
     "Touchez une phrase\n→ Traduit instantanément\n→ Toutes langues"),
    ("03", "Conversation\nA/B Bilingue",
     "2 personnes · 2 langues\nTraduction + audio auto",
     "Personne A parle\n→ Personne B entend\n→ Dans sa langue"),
    ("04", "Carte JOJ\nInteractive",
     "8 sites · 3 zones\nOpenStreetMap + Google Maps",
     "Choisissez un site\n→ Carte interactive\n→ Navigation GPS"),
    ("05", "Assistant JOJ\nIA Multilingue",
     "22 intentions · Transport\nSIM · Urgences · Culture",
     "Posez une question\n→ Réponse en votre langue\n→ Liens d'action"),
]

col_positions = [0.28, 2.87, 5.46, 8.05, 10.64]
for i, (num, title, desc, flow) in enumerate(features):
    lx = col_positions[i]
    # Carte fond
    add_rect(s4, lx, 1.35, 2.4, 5.8, GRAY_LIGHT)
    # Header coloré
    add_rect(s4, lx, 1.35, 2.4, 0.65, GREEN)
    add_text(s4, num, lx+0.08, 1.38, 0.6, 0.55,
             font_size=20, bold=True, color=RGBColor(0xA8, 0xE6, 0xC3))
    add_text(s4, title, lx+0.1, 1.38, 2.15, 0.6,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    # Description
    add_text(s4, desc, lx+0.1, 2.1, 2.2, 1.2,
             font_size=12, color=DARK_TEXT)
    # Séparateur
    add_rect(s4, lx+0.1, 3.35, 2.2, 0.04, GREEN)
    # Flow utilisateur
    add_text(s4, "Parcours :", lx+0.1, 3.45, 2.2, 0.3,
             font_size=11, bold=True, color=GREEN_DARK)
    add_text(s4, flow, lx+0.1, 3.78, 2.2, 2.8,
             font_size=11, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MARCHÉ / CIBLE
# ═══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, 13.33, 7.5, WHITE)
green_header_bar(s5, "Marché / Cible", "De Dakar 2026 à l'Afrique entière")
footer_band(s5)
slide_number(s5, 5)

# Colonne gauche — marché immédiat
add_rect(s5, 0.3, 1.4, 5.8, 5.7, GREEN_DEEP)
add_text(s5, "MARCHÉ IMMÉDIAT", 0.5, 1.5, 5.4, 0.45,
         font_size=15, bold=True, color=GREEN)
add_text(s5, "JOJ Dakar 2026", 0.5, 1.98, 5.4, 0.55,
         font_size=22, bold=True, color=WHITE)

imm_stats = [
    ("4 000", "Athlètes · 206 nations"),
    ("10 000", "Bénévoles & staff JOJ"),
    ("50 000+", "Visiteurs internationaux"),
    ("3 000", "Médias & journalistes"),
]
for i, (num, label) in enumerate(imm_stats):
    ly = 2.65 + i * 0.85
    add_rect(s5, 0.45, ly, 5.55, 0.72, RGBColor(0x0F, 0x3A, 0x22))
    add_text(s5, num, 0.55, ly+0.05, 1.8, 0.6,
             font_size=24, bold=True, color=GREEN)
    add_text(s5, label, 2.4, ly+0.18, 3.4, 0.4,
             font_size=14, color=WHITE)

add_rect(s5, 0.45, 6.08, 5.55, 0.6, GREEN)
add_text(s5, "= ~67 000 utilisateurs potentiels", 0.55, 6.13, 5.3, 0.5,
         font_size=15, bold=True, color=GREEN_DEEP, align=PP_ALIGN.CENTER)

# Colonne droite — extension
add_rect(s5, 6.5, 1.4, 6.53, 5.7, GRAY_LIGHT)
add_text(s5, "EXTENSION · AFRIQUE & MONDE", 6.65, 1.5, 6.2, 0.4,
         font_size=13, bold=True, color=GREEN_DARK)

events = [
    ("2025", "CAN Maroc", "Arabe · Français · Anglais"),
    ("2026", "JOJ Dakar", "13 langues dont Wolof ← Nous"),
    ("2027", "FIBA AfroBasket", "Swahili · Anglais · Français"),
    ("2028", "LA Olympics", "Espagnol · Anglais · Chinois"),
    ("2028+", "Sommets UA / Forums", "Toutes langues africaines"),
]

for i, (year, event, langs) in enumerate(events):
    ly = 2.0 + i * 0.95
    is_current = "Nous" in langs
    bg = GREEN if is_current else WHITE
    fg = WHITE if is_current else DARK_TEXT
    add_rect(s5, 6.6, ly, 6.1, 0.82, bg)
    add_text(s5, year, 6.68, ly+0.04, 0.8, 0.35,
             font_size=18, bold=True,
             color=GREEN if not is_current else WHITE)
    add_text(s5, event, 7.55, ly+0.04, 3.5, 0.35,
             font_size=15, bold=True, color=fg)
    add_text(s5, langs, 7.55, ly+0.42, 4.8, 0.3,
             font_size=12, color=fg if is_current else GRAY_MED)

add_text(s5, "Marché total : 150 M€+ d'événements int. en Afrique (2024–2032)",
         6.6, 6.08, 6.15, 0.55,
         font_size=13, bold=True, color=GREEN_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, 13.33, 7.5, WHITE)
green_header_bar(s6, "Business Model", "Un modèle hybride B2B + Affiliation")
footer_band(s6)
slide_number(s6, 6)

# 3 colonnes temporelles
columns = [
    ("COURT TERME\nJOJ 2026", GREEN_DEEP, [
        ("Licence officielle B2B", "CIO / CNOSS", "50K – 200K €"),
        ("Sponsoring in-app", "Orange · Yango", "10K – 50K €"),
    ]),
    ("MOYEN TERME\nPost-JOJ", GREEN_DARK, [
        ("Affiliation transport", "Commission Yango", "2 – 5 € / course"),
        ("Listings premium", "Hôtels · Restos", "50 – 200 € / mois"),
    ]),
    ("LONG TERME\nÉvénements", GREEN, [
        ("Licences événements", "CAN · FIBA · Sommets", "Récurrent"),
        ("API as a Service", "Trad. événements tiers", "SaaS"),
    ]),
]

for ci, (period, col, items) in enumerate(columns):
    lx = 0.4 + ci * 4.3
    add_rect(s6, lx, 1.38, 4.0, 0.65, col)
    add_text(s6, period, lx+0.1, 1.4, 3.8, 0.6,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for ri, (source, detail, amount) in enumerate(items):
        ly = 2.15 + ri * 1.5
        add_rect(s6, lx, ly, 4.0, 1.35, GRAY_LIGHT)
        add_rect(s6, lx, ly, 0.15, 1.35, col)
        add_text(s6, source, lx+0.25, ly+0.08, 3.65, 0.38,
                 font_size=15, bold=True, color=DARK_TEXT)
        add_text(s6, detail, lx+0.25, ly+0.48, 3.65, 0.3,
                 font_size=12, color=GRAY_MED)
        add_text(s6, amount, lx+0.25, ly+0.82, 3.65, 0.38,
                 font_size=16, bold=True, color=col)

# Structure de coûts (bas)
add_rect(s6, 0.4, 5.35, 12.53, 1.55, GREEN_LIGHT)
add_text(s6, "STRUCTURE DE COÛTS", 0.6, 5.42, 4, 0.35,
         font_size=12, bold=True, color=GREEN_DARK)

costs = [
    ("API DeepL", "~200-800 €/mois JOJ"),
    ("Hébergement Vercel", "~100-300 €/mois"),
    ("Développement", "Équipe fondatrice"),
    ("Marketing", "Partenariats JOJ"),
]
for i, (label, val) in enumerate(costs):
    lx = 0.6 + i * 3.15
    add_text(s6, label, lx, 5.82, 2.8, 0.3,
             font_size=12, bold=True, color=DARK_TEXT)
    add_text(s6, val, lx, 6.12, 2.8, 0.55,
             font_size=13, color=GREEN_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — COLLABORATION ORANGE ⭐
# ═══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, 13.33, 7.5, WHITE)

# Header Orange (couleurs Orange Sénégal)
add_rect(s7, 0, 0, 13.33, 1.25, RGBColor(0xFF, 0x7A, 0x1A))
add_text(s7, "Collaboration avec Orange", 0.4, 0.08, 10, 0.6,
         font_size=30, bold=True, color=WHITE)
add_text(s7, "4 intégrations concrètes · Valeur mutuelle · Immédiatement faisables",
         0.4, 0.72, 12, 0.42,
         font_size=15, color=RGBColor(0xFF, 0xE5, 0xCC))

footer_band(s7)
slide_number(s7, 7)

collabs = [
    (
        "01",
        "Bundle SIM Touriste × DÉGG",
        "La plus impactante",
        "Chaque SIM Touriste Orange vendue à l'AIBD\ncontient un QR code DÉGG + forfait data activé.",
        "Orange : consommation data + visibilité 50K visiteurs\nDÉGG : acquisition automatique de tous les arrivants",
        "Distribution"
    ),
    (
        "02",
        "DÉGG dans Orange Max IT",
        "Distribution massive",
        "DÉGG intégré comme mini-app dans l'écosystème\nMax IT · Push notifications JOJ multilingues.",
        "Orange : Max IT = super-app des JOJ\nDÉGG : accès 10M+ abonnés Orange Sénégal",
        "Max IT"
    ),
    (
        "03",
        "Réseau Officiel · Co-branding",
        "Légitimité immédiate",
        "« DÉGG propulsé par Orange Sénégal »\nOrange garantit la couverture réseau sur les 8 sites.",
        "Orange : co-branding JOJ + 206 nations\nDÉGG : badge officiel + fiabilité réseau",
        "Réseau"
    ),
    (
        "04",
        "Orange Money + SMS API",
        "Tech & Paiement",
        "Achat billets JOJ via Orange Money depuis DÉGG\nAlertes SMS multilingues via API Orange SMS.",
        "Orange : transactions + usage SMS API\nDÉGG : paiement natif + alertes push",
        "FinTech"
    ),
]

positions_7 = [
    (0.3, 1.38),
    (6.7, 1.38),
    (0.3, 4.3),
    (6.7, 4.3),
]

for i, ((lx, ly), (num, title, tag, desc, value, badge_txt)) in enumerate(
        zip(positions_7, collabs)):
    add_rect(s7, lx, ly, 6.2, 2.75, GRAY_LIGHT)
    add_rect(s7, lx, ly, 6.2, 0.55, ORANGE_ACC)

    # Numéro
    add_text(s7, num, lx+0.1, ly+0.06, 0.55, 0.45,
             font_size=20, bold=True, color=WHITE)
    # Titre
    add_text(s7, title, lx+0.7, ly+0.06, 4.5, 0.45,
             font_size=16, bold=True, color=WHITE)
    # Badge tag
    add_badge(s7, badge_txt, lx+5.15, ly+0.1, 0.9, 0.32,
              bg=GREEN, fg=WHITE, font_size=10)

    # Description
    add_text(s7, desc, lx+0.12, ly+0.65, 5.9, 1.05,
             font_size=13, color=DARK_TEXT)

    # Valeur pour chaque partie
    add_rect(s7, lx+0.12, ly+1.75, 5.9, 0.85, WHITE)
    add_text(s7, value, lx+0.2, ly+1.82, 5.7, 0.75,
             font_size=12, color=GREEN_DARK)

# Message final
add_rect(s7, 0.3, 7.05, 12.73, 0.35, GREEN_DEEP)
add_text(s7, "L'opportunité pour Orange, c'est MAINTENANT — avant les JOJ, pas après.",
         0.5, 7.07, 12.3, 0.28,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─── Sauvegarde ──────────────────────────────────────────────────────────────
output_path = r"c:\Users\p\Downloads\DEGG_Pitch_Orange_Hackathon.pptx"
prs.save(output_path)
print(f"OK Fichier genere : {output_path}")
print(f"   Taille : {os.path.getsize(output_path) / 1024:.1f} Ko")
